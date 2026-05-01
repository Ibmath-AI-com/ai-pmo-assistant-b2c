from __future__ import annotations

import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def generate_pptx(input_data: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides_data = input_data.get("slide_data", [])

    if not slides_data:
        slides_data = [{"title": input_data.get("title", "Presentation"), "content": input_data.get("content", ""), "layout": "title"}]

    for i, slide_info in enumerate(slides_data):
        layout_name = slide_info.get("layout", "content")

        if layout_name == "title" or i == 0:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_info.get("title", "")
            placeholders = [ph for ph in slide.placeholders if ph.placeholder_format.idx == 1]
            if placeholders:
                placeholders[0].text = slide_info.get("subtitle", slide_info.get("content", ""))
        else:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_info.get("title", "")

            content = slide_info.get("content", "")
            bullets = slide_info.get("bullets", [])

            body_placeholders = [ph for ph in slide.placeholders if ph.placeholder_format.idx == 1]
            if body_placeholders:
                tf = body_placeholders[0].text_frame
                tf.clear()

                if bullets:
                    for j, bullet in enumerate(bullets):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
                elif content:
                    tf.paragraphs[0].text = content

    prs.save(output_path)
    return output_path


def create_temp_pptx(input_data: dict) -> str:
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.close()
    generate_pptx(input_data, tmp.name)
    return tmp.name
